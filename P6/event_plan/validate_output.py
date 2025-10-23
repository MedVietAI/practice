#!/usr/bin/env python3
"""
Validation script for Vietnam 80th Anniversary Event Planning System
Comprehensive quality assurance and impact assessment
"""

import os
import json
from pptx import Presentation
from config import OUTPUT_DIR, DIST_DIR, ASSETS_DIR

def validate_content_quality():
    """Validate the quality and comprehensiveness of generated content"""
    print("🔍 Validating Content Quality...")
    
    content_file = os.path.join(OUTPUT_DIR, 'content.json')
    if not os.path.exists(content_file):
        print("❌ Content file not found")
        return False
    
    with open(content_file, 'r', encoding='utf-8') as f:
        content = json.load(f)
    
    # Check required sections
    required_sections = ['title', 'objectives', 'scope', 'stakeholders', 'timeline', 
                       'program', 'communications', 'logistics', 'budget', 'risk', 
                       'safety', 'sustainability', 'raci', 'kpi', 'approvals']
    
    missing_sections = [section for section in required_sections if section not in content]
    if missing_sections:
        print(f"❌ Missing sections: {missing_sections}")
        return False
    
    # Check content richness
    objectives_count = len(content.get('objectives', []))
    stakeholders_count = len(content.get('stakeholders', []))
    timeline_phases = len(content.get('timeline', []))
    programs_count = len(content.get('program', []))
    
    print(f"✅ Content Quality Metrics:")
    print(f"   📊 Objectives: {objectives_count} (Target: ≥4)")
    print(f"   📊 Stakeholders: {stakeholders_count} (Target: ≥6)")
    print(f"   📊 Timeline Phases: {timeline_phases} (Target: 4)")
    print(f"   📊 Programs: {programs_count} (Target: ≥3)")
    
    # Check patriotic language
    patriotic_keywords = ['tổ quốc', 'dân tộc', 'đoàn kết', 'yêu nước', 'tự hào', 'khát vọng']
    content_text = json.dumps(content, ensure_ascii=False).lower()
    patriotic_score = sum(1 for keyword in patriotic_keywords if keyword in content_text)
    
    print(f"   🇻🇳 Patriotic Language Score: {patriotic_score}/6 keywords found")
    
    return objectives_count >= 4 and stakeholders_count >= 6 and timeline_phases == 4

def validate_powerpoint_quality():
    """Validate PowerPoint presentation quality"""
    print("\n🎨 Validating PowerPoint Quality...")
    
    pptx_file = os.path.join(DIST_DIR, 'plan_80nam_A90.pptx')
    if not os.path.exists(pptx_file):
        print("❌ PowerPoint file not found")
        return False
    
    try:
        prs = Presentation(pptx_file)
        slide_count = len(prs.slides)
        
        print(f"✅ PowerPoint Quality Metrics:")
        print(f"   📊 Total Slides: {slide_count} (Target: ≥15)")
        print(f"   📊 File Size: {os.path.getsize(pptx_file) / 1024 / 1024:.1f} MB")
        
        # Check slide content
        title_slides = 0
        content_slides = 0
        
        for slide in prs.slides:
            if slide.shapes:
                # Count text boxes
                text_boxes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame')]
                if len(text_boxes) >= 2:  # Title + content
                    content_slides += 1
                elif len(text_boxes) == 1:  # Title slide
                    title_slides += 1
        
        print(f"   📊 Title Slides: {title_slides}")
        print(f"   📊 Content Slides: {content_slides}")
        
        return slide_count >= 15 and content_slides >= 10
        
    except Exception as e:
        print(f"❌ Error reading PowerPoint: {e}")
        return False

def validate_image_assets():
    """Validate crawled image assets"""
    print("\n🖼️ Validating Image Assets...")
    
    images_file = os.path.join(ASSETS_DIR, 'images.json')
    if not os.path.exists(images_file):
        print("❌ Images file not found")
        return False
    
    with open(images_file, 'r', encoding='utf-8') as f:
        images = json.load(f)
    
    total_images = len(images)
    patriotic_images = 0
    
    patriotic_keywords = ['quốc khánh', 'cách mạng', 'độc lập', 'tự do', 'tổ quốc', 'việt nam', 'hồ chí minh', 'ba đình']
    
    for img in images:
        alt_text = img.get('alt', '').lower()
        title_text = img.get('title', '').lower()
        if any(keyword in alt_text or keyword in title_text for keyword in patriotic_keywords):
            patriotic_images += 1
    
    print(f"✅ Image Assets Metrics:")
    print(f"   📊 Total Images: {total_images} (Target: ≥500)")
    print(f"   📊 Patriotic Images: {patriotic_images} (Target: ≥50)")
    print(f"   📊 Patriotic Ratio: {patriotic_images/total_images*100:.1f}%")
    
    return total_images >= 500 and patriotic_images >= 50

def validate_storytelling_impact():
    """Validate storytelling and emotional impact"""
    print("\n📖 Validating Storytelling Impact...")
    
    content_file = os.path.join(OUTPUT_DIR, 'content.json')
    with open(content_file, 'r', encoding='utf-8') as f:
        content = json.load(f)
    
    # Check narrative flow
    timeline = content.get('timeline', [])
    narrative_flow = len([phase for phase in timeline if 'giai đoạn' in phase.get('phase', '').lower()])
    
    # Check emotional language
    emotional_keywords = ['vẻ vang', 'hào hùng', 'tự hào', 'khát vọng', 'đoàn kết', 'yêu nước']
    content_text = json.dumps(content, ensure_ascii=False).lower()
    emotional_score = sum(1 for keyword in emotional_keywords if keyword in content_text)
    
    # Check historical significance
    historical_keywords = ['cách mạng', 'độc lập', 'tự do', 'tổ quốc', 'dân tộc']
    historical_score = sum(1 for keyword in historical_keywords if keyword in content_text)
    
    print(f"✅ Storytelling Impact Metrics:")
    print(f"   📊 Narrative Flow: {narrative_flow}/4 phases with clear progression")
    print(f"   📊 Emotional Language: {emotional_score}/6 keywords found")
    print(f"   📊 Historical Significance: {historical_score}/5 keywords found")
    
    return narrative_flow == 4 and emotional_score >= 4 and historical_score >= 4

def generate_impact_report():
    """Generate comprehensive impact report"""
    print("\n📊 Generating Comprehensive Impact Report...")
    
    # Run all validations
    content_quality = validate_content_quality()
    pptx_quality = validate_powerpoint_quality()
    image_assets = validate_image_assets()
    storytelling = validate_storytelling_impact()
    
    # Calculate overall score
    total_score = sum([content_quality, pptx_quality, image_assets, storytelling])
    max_score = 4
    
    print(f"\n🎯 COMPREHENSIVE IMPACT ASSESSMENT")
    print(f"{'='*50}")
    print(f"📊 Overall Quality Score: {total_score}/{max_score} ({total_score/max_score*100:.1f}%)")
    print(f"")
    print(f"✅ Content Quality: {'PASS' if content_quality else 'FAIL'}")
    print(f"✅ PowerPoint Quality: {'PASS' if pptx_quality else 'FAIL'}")
    print(f"✅ Image Assets: {'PASS' if image_assets else 'FAIL'}")
    print(f"✅ Storytelling Impact: {'PASS' if storytelling else 'FAIL'}")
    print(f"")
    
    if total_score == max_score:
        print("🎉 EXCELLENT! All quality metrics passed!")
        print("🇻🇳 High patriotism and storytelling impact achieved!")
        print("📈 Comprehensive and meaningful product delivered!")
    elif total_score >= 3:
        print("✅ GOOD! Most quality metrics passed!")
        print("📊 Minor improvements possible")
    else:
        print("⚠️  NEEDS IMPROVEMENT! Some quality metrics failed!")
        print("🔧 Consider regenerating content or fixing issues")
    
    return total_score == max_score

if __name__ == "__main__":
    print("🇻🇳 Vietnam 80th Anniversary Event Planning - Quality Validation")
    print("="*70)
    
    success = generate_impact_report()
    
    if success:
        print("\n🎊 VALIDATION SUCCESSFUL!")
        print("The product meets all quality standards for:")
        print("• Comprehensive content coverage")
        print("• Professional presentation quality")
        print("• Rich patriotic imagery")
        print("• High storytelling impact")
        print("• Meaningful and impactful delivery")
    else:
        print("\n⚠️  VALIDATION INCOMPLETE!")
        print("Please review the failed metrics and regenerate if needed.")

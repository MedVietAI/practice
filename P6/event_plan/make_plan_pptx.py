#!/usr/bin/env python3
"""
PowerPoint generator for Vietnam 80th Anniversary event plan
Creates professional presentation with comprehensive content
"""

import os
import json
import random
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.dml import MSO_FILL
from PIL import Image
from config import OUTPUT_DIR, DIST_DIR, EVENT_TITLE, EVENT_FRAMEWORK, ASSETS_DIR

# Slide dimensions (16:9)
W, H = Inches(13.33), Inches(7.5)

# Color scheme (Vietnamese flag colors)
THEME_BG = RGBColor(230, 0, 0)      # Red
ACCENT = RGBColor(255, 210, 0)      # Gold
TEXT = RGBColor(20, 20, 20)         # Dark text
WHITE = RGBColor(255, 255, 255)     # White
LIGHT_BG = RGBColor(248, 248, 248)  # Light background

class PowerPointGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.slide_w = self.prs.slide_width
        self.slide_h = self.prs.slide_height
        self.images_data = self.load_images_data()
    
    def load_images_data(self):
        """Load crawled images data"""
        try:
            images_file = os.path.join(ASSETS_DIR, 'images.json')
            if os.path.exists(images_file):
                with open(images_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
        except Exception as e:
            print(f"⚠️  Could not load images: {e}")
        return []
    
    def get_patriotic_image(self, keywords=None):
        """Get a patriotic image from crawled data"""
        if not self.images_data:
            return None
        
        # Filter images by patriotic keywords
        patriotic_keywords = ['quốc khánh', 'cách mạng', 'độc lập', 'tự do', 'tổ quốc', 'việt nam', 'hồ chí minh', 'ba đình']
        
        for img in self.images_data:
            alt_text = img.get('alt', '').lower()
            title_text = img.get('title', '').lower()
            
            if any(keyword in alt_text or keyword in title_text for keyword in patriotic_keywords):
                return img
        
        # Fallback to any image
        return self.images_data[0] if self.images_data else None
    
    def download_image(self, image_url, max_size=(800, 600)):
        """Download and resize image for PowerPoint"""
        try:
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            
            # Open image with PIL
            img = Image.open(BytesIO(response.content))
            
            # Convert to RGB if necessary
            if img.mode in ('RGBA', 'LA', 'P'):
                img = img.convert('RGB')
            
            # Resize image while maintaining aspect ratio
            img.thumbnail(max_size, Image.Resampling.LANCZOS)
            
            # Save to temporary file
            temp_path = f"/tmp/temp_image_{random.randint(1000, 9999)}.jpg"
            img.save(temp_path, 'JPEG', quality=85)
            
            return temp_path
        except Exception as e:
            print(f"⚠️  Could not download image {image_url}: {e}")
            return None
    
    def add_image_to_slide(self, slide, image_path, x, y, width, height, opacity=0.3):
        """Add image to slide with specified dimensions and opacity"""
        try:
            if os.path.exists(image_path):
                # Add image
                pic = slide.shapes.add_picture(image_path, x, y, width, height)
                
                # Set transparency if opacity < 1
                if opacity < 1:
                    pic.fill.solid()
                    pic.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    # Note: python-pptx has limited transparency support
                
                return pic
        except Exception as e:
            print(f"⚠️  Could not add image to slide: {e}")
        return None
        
    def add_title_slide(self, title, subtitle=""):
        """Create title slide with Vietnamese theme and patriotic imagery"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background with gradient effect
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.slide_w, self.slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = THEME_BG
        bg.line.fill.background()
        
        # Add patriotic elements
        self.add_patriotic_elements(slide)
        
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(2.5))
        title_frame = title_box.text_frame
        title_frame.clear()
        
        # Main title
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(48)
        title_p.font.bold = True
        title_p.font.color.rgb = WHITE
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_p = title_frame.add_paragraph()
            subtitle_p.text = subtitle
            subtitle_p.font.size = Pt(24)
            subtitle_p.font.color.rgb = WHITE
            subtitle_p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_patriotic_elements(self, slide):
        """Add patriotic visual elements to slides"""
        # Add golden star in corner
        star = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.STAR_5_POINT, 
                                    Inches(0.5), Inches(0.5), Inches(1), Inches(1))
        star.fill.solid()
        star.fill.fore_color.rgb = ACCENT
        star.line.fill.background()
        
        # Add golden border
        border = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 
                                      Inches(0.1), Inches(0.1), 
                                      self.slide_w - Inches(0.2), self.slide_h - Inches(0.2))
        border.fill.background()
        border.line.color.rgb = ACCENT
        border.line.width = Pt(3)
    
    def add_patriotic_hero_slide(self, title, content, image_url=None):
        """Create a patriotic hero slide with background image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.slide_w, self.slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = THEME_BG
        bg.line.fill.background()
        
        # Add patriotic elements
        self.add_patriotic_elements(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.5), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = WHITE
        title_p.alignment = PP_ALIGN.CENTER
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.5), Inches(4))
        content_frame = content_box.text_frame
        content_frame.clear()
        
        for i, item in enumerate(content):
            p = content_frame.add_paragraph() if i else content_frame.paragraphs[0]
            p.text = f"★ {item}"
            p.font.size = Pt(24)
            p.font.color.rgb = WHITE
            p.font.bold = True
        
        return slide
    
    def add_storyful_slide(self, title, content, image_keywords=None, contrast_theme=None):
        """Create a storyful slide with historical contrast and imagery"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Get appropriate image
        image_data = self.get_patriotic_image(image_keywords)
        image_path = None
        
        if image_data and image_data.get('url'):
            image_path = self.download_image(image_data['url'])
        
        # Background
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.slide_w, self.slide_h)
        if contrast_theme == "historical":
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(50, 50, 50)  # Dark for historical
        else:
            bg.fill.solid()
            bg.fill.fore_color.rgb = LIGHT_BG
        
        # Add background image if available
        if image_path:
            self.add_image_to_slide(slide, image_path, Inches(0.5), Inches(1), 
                                  Inches(12.5), Inches(5.5), opacity=0.2)
        
        # Add patriotic elements
        self.add_patriotic_elements(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.8), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = THEME_BG if contrast_theme != "historical" else WHITE
        title_p.alignment = PP_ALIGN.CENTER
        
        # Content with story elements
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.8), Inches(5))
        content_frame = content_box.text_frame
        content_frame.clear()
        
        for i, item in enumerate(content):
            p = content_frame.add_paragraph() if i else content_frame.paragraphs[0]
            
            # Add story elements based on theme
            if contrast_theme == "historical":
                p.text = f"⚔️ {item}"
            elif contrast_theme == "modern":
                p.text = f"🏛️ {item}"
            else:
                p.text = f"★ {item}"
            
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT if contrast_theme != "historical" else WHITE
            p.font.bold = True
        
        return slide
    
    def add_historical_contrast_slide(self, title, historical_content, modern_content):
        """Create a slide showing historical sacrifice vs modern peace"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background with gradient effect
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.slide_w, self.slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(240, 240, 240)
        bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = THEME_BG
        title_p.alignment = PP_ALIGN.CENTER
        
        # Historical side (left)
        hist_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        hist_frame = hist_box.text_frame
        hist_frame.clear()
        
        # Historical background
        hist_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 
                                       Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        hist_bg.fill.solid()
        hist_bg.fill.fore_color.rgb = RGBColor(100, 50, 50)  # Dark red for sacrifice
        hist_bg.line.fill.background()
        
        hist_title = hist_frame.paragraphs[0]
        hist_title.text = "⚔️ THỜI KỲ HY SINH (1945-1975)"
        hist_title.font.size = Pt(18)
        hist_title.font.bold = True
        hist_title.font.color.rgb = WHITE
        
        for i, item in enumerate(historical_content):
            p = hist_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
        
        # Modern side (right)
        mod_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5))
        mod_frame = mod_box.text_frame
        mod_frame.clear()
        
        # Modern background
        mod_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 
                                      Inches(6.8), Inches(1.5), Inches(6), Inches(5))
        mod_bg.fill.solid()
        mod_bg.fill.fore_color.rgb = RGBColor(50, 100, 50)  # Green for peace/prosperity
        mod_bg.line.fill.background()
        
        mod_title = mod_frame.paragraphs[0]
        mod_title.text = "🏛️ THỜI KỲ HÒA BÌNH (1975-2025)"
        mod_title.font.size = Pt(18)
        mod_title.font.bold = True
        mod_title.font.color.rgb = WHITE
        
        for i, item in enumerate(modern_content):
            p = mod_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
        
        return slide
    
    def add_bullets_slide(self, title, bullets, two_col=False):
        """Create bullet points slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Light background
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.slide_w, self.slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.8), Inches(1.0))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = THEME_BG
        title_p.alignment = PP_ALIGN.CENTER
        
        if two_col:
            # Two column layout
            x1, y, w, h = Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.6)
            x2 = Inches(6.6)
            cols = [(x1, y, w, h), (x2, y, w, h)]
            chunks = [bullets[:len(bullets)//2+1], bullets[len(bullets)//2+1:]]
            
            for (x, y, w, h), bullet_list in zip(cols, chunks):
                text_box = slide.shapes.add_textbox(x, y, w, h)
                text_frame = text_box.text_frame
                text_frame.clear()
                
                for i, bullet in enumerate(bullet_list):
                    p = text_frame.add_paragraph() if i else text_frame.paragraphs[0]
                    p.text = str(bullet)
                    p.level = 0
                    p.font.size = Pt(20)
                    p.font.color.rgb = TEXT
        else:
            # Single column layout
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.8))
            text_frame = text_box.text_frame
            text_frame.clear()
            
            for i, bullet in enumerate(bullets):
                p = text_frame.add_paragraph() if i else text_frame.paragraphs[0]
                p.text = str(bullet)
                p.level = 0
                p.font.size = Pt(22)
                p.font.color.rgb = TEXT
        
        return slide
    
    def add_timeline_slide(self, title, timeline_data):
        """Create timeline slide with visual elements"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.slide_w, self.slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.8), Inches(1.0))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = THEME_BG
        title_p.alignment = PP_ALIGN.CENTER
        
        # Timeline content
        y_start = Inches(1.8)
        for i, phase in enumerate(timeline_data):
            # Phase box
            phase_box = slide.shapes.add_textbox(Inches(0.8), y_start + i * Inches(1.2), Inches(11.8), Inches(1.0))
            phase_frame = phase_box.text_frame
            phase_frame.clear()
            
            # Phase title
            phase_p = phase_frame.paragraphs[0]
            phase_p.text = f"{phase.get('phase', '')} — {phase.get('window', '')}"
            phase_p.font.size = Pt(24)
            phase_p.font.bold = True
            phase_p.font.color.rgb = THEME_BG
            
            # Workstreams
            for j, workstream in enumerate(phase.get('workstreams', [])):
                ws_p = phase_frame.add_paragraph()
                ws_p.text = f"  • {workstream}"
                ws_p.font.size = Pt(18)
                ws_p.font.color.rgb = TEXT
                ws_p.level = 1
        
        return slide
    
    def add_budget_slide(self, title, budget_data):
        """Create budget slide with table-like layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.slide_w, self.slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.8), Inches(1.0))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = THEME_BG
        title_p.alignment = PP_ALIGN.CENTER
        
        # Budget items
        y_start = Inches(1.8)
        for i, item in enumerate(budget_data):
            item_box = slide.shapes.add_textbox(Inches(0.8), y_start + i * Inches(1.0), Inches(11.8), Inches(0.8))
            item_frame = item_box.text_frame
            item_frame.clear()
            
            item_p = item_frame.paragraphs[0]
            item_p.text = f"{item.get('category', '')}: {item.get('planned', '')} ({item.get('notes', '')})"
            item_p.font.size = Pt(20)
            item_p.font.color.rgb = TEXT
        
        return slide
    
    def add_risk_matrix_slide(self, title, risk_data):
        """Create risk assessment slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.slide_w, self.slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.8), Inches(1.0))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = THEME_BG
        title_p.alignment = PP_ALIGN.CENTER
        
        # Risk items
        y_start = Inches(1.8)
        for i, risk in enumerate(risk_data):
            risk_box = slide.shapes.add_textbox(Inches(0.8), y_start + i * Inches(1.2), Inches(11.8), Inches(1.0))
            risk_frame = risk_box.text_frame
            risk_frame.clear()
            
            risk_p = risk_frame.paragraphs[0]
            risk_p.text = f"{risk.get('risk', '')} — Tác động: {risk.get('impact', '')} — Xác suất: {risk.get('likelihood', '')}"
            risk_p.font.size = Pt(20)
            risk_p.font.bold = True
            risk_p.font.color.rgb = THEME_BG
            
            mitigation_p = risk_frame.add_paragraph()
            mitigation_p.text = f"Ứng phó: {risk.get('mitigation', '')}"
            mitigation_p.font.size = Pt(18)
            mitigation_p.font.color.rgb = TEXT
            mitigation_p.level = 1
        
        return slide
    
    def generate_presentation(self, data):
        """Generate complete presentation from data"""
        print("📊 Generating PowerPoint presentation...")
        
        # 1) Title slide
        self.add_title_slide(
            data.get('title', EVENT_TITLE), 
            f'Khung {EVENT_FRAMEWORK} — T-60 → T+30'
        )
        
        # 2) Historical Context & Objectives
        historical_sacrifices = [
            "Cách mạng Tháng Tám 1945 - Đánh đuổi thực dân Pháp",
            "Kháng chiến chống Mỹ cứu nước (1954-1975)",
            "Hy sinh của hàng triệu đồng bào vì độc lập tự do",
            "Tinh thần bất khuất của dân tộc Việt Nam"
        ]
        
        modern_achievements = [
            "Hòa bình, thống nhất đất nước từ 1975",
            "Đổi mới và phát triển kinh tế từ 1986",
            "Hội nhập quốc tế và phát triển bền vững",
            "Khát vọng trở thành nước phát triển vào 2045"
        ]
        
        self.add_historical_contrast_slide(
            "Từ Hy Sinh Đến Hòa Bình - 80 Năm Độc Lập",
            historical_sacrifices,
            modern_achievements
        )
        
        # 3) Objectives with patriotic imagery
        self.add_storyful_slide('Mục tiêu Cao Cả', data.get('objectives', []), 
                              image_keywords=['quốc khánh', 'cách mạng'], contrast_theme="modern")
        
        # 4) Scope
        self.add_bullets_slide('Phạm vi', data.get('scope', []))
        
        # 5) Stakeholders
        self.add_bullets_slide('Các bên liên quan', data.get('stakeholders', []), two_col=True)
        
        # 6) Timeline with historical context
        self.add_timeline_slide('Dòng thời gian (A90)', data.get('timeline', []))
        
        # 7) Key programs with imagery
        programs = [f"{p.get('name', '')}: {p.get('goal', '')} (Chủ trì: {p.get('owner', '')})" 
                   for p in data.get('program', [])]
        self.add_storyful_slide('Chương trình trọng điểm', programs, 
                              image_keywords=['lễ kỷ niệm', 'triển lãm'], contrast_theme="modern")
        
        # 6) Communications
        communications = [f"{c.get('channel', '')}: {c.get('message', '')} — Mốc: {c.get('key_dates', '')}" 
                        for c in data.get('communications', [])]
        self.add_bullets_slide('Kế hoạch truyền thông', communications, two_col=True)
        
        # 7) Logistics
        logistics = [f"{l.get('item', '')}: {l.get('details', '')} — Hạn: {l.get('deadline', '')}" 
                   for l in data.get('logistics', [])]
        self.add_bullets_slide('Hậu cần', logistics)
        
        # 8) Budget
        self.add_budget_slide('Ngân sách (dự trù)', data.get('budget', []))
        
        # 9) Risk assessment
        self.add_risk_matrix_slide('Rủi ro chính & Ứng phó', data.get('risk', []))
        
        # 10) Safety
        safety = [f"{s.get('topic', '')}: {s.get('actions', '')}" for s in data.get('safety', [])]
        self.add_bullets_slide('An toàn', safety)
        
        # 11) Sustainability
        sustainability = [f"{s.get('topic', '')}: {s.get('actions', '')}" for s in data.get('sustainability', [])]
        self.add_bullets_slide('Bền vững', sustainability)
        
        # 12) RACI Matrix
        raci = [f"{r.get('task', '')}: R={r.get('R', '')}, A={r.get('A', '')}, C={r.get('C', '')}, I={r.get('I', '')}" 
               for r in data.get('raci', [])]
        self.add_bullets_slide('Ma trận RACI', raci)
        
        # 13) KPIs
        kpis = [f"{k.get('metric', '')}: {k.get('definition', '')} — Mục tiêu: {k.get('target', '')}" 
               for k in data.get('kpi', [])]
        self.add_bullets_slide('Chỉ số đo lường (KPI)', kpis)
        
        # 14) Approvals
        approvals = [f"Cổng {a.get('gate', '')}: {a.get('criteria', '')} — Chủ trì: {a.get('owner', '')}" 
                    for a in data.get('approvals', [])]
        self.add_bullets_slide('Cổng phê duyệt', approvals)
        
        # 15) Journey from War to Peace
        war_to_peace_content = [
            "Từ những năm tháng khói lửa đến hòa bình thịnh vượng",
            "Từ hy sinh anh dũng đến thành tựu vẻ vang",
            "Từ đất nước bị chia cắt đến thống nhất toàn vẹn",
            "Từ nghèo đói đến phát triển, từ lạc hậu đến hiện đại"
        ]
        self.add_storyful_slide('Hành Trình 80 Năm', war_to_peace_content, 
                              image_keywords=['hòa bình', 'phát triển'], contrast_theme="modern")
        
        # 16) Patriotic closing slide
        patriotic_quotes = [
            "Việt Nam muôn năm!",
            "Độc lập - Tự do - Hạnh phúc",
            "Đoàn kết toàn dân tộc",
            "Khát vọng Rồng bay"
        ]
        self.add_patriotic_hero_slide('Tri ân — Đoàn kết — Khát vọng', patriotic_quotes)
        
        print(f"✅ Generated {len(self.prs.slides)} slides")
    
    def save_presentation(self, filename="plan_80nam_A90.pptx"):
        """Save presentation to file"""
        os.makedirs(DIST_DIR, exist_ok=True)
        filepath = os.path.join(DIST_DIR, filename)
        self.prs.save(filepath)
        print(f"✅ Saved presentation to {filepath}")
        return filepath

def main():
    """Main function to generate PowerPoint presentation"""
    print("🎯 Starting PowerPoint generation for Vietnam 80th Anniversary...")
    
    # Load content
    content_file = os.path.join(OUTPUT_DIR, 'content.json')
    if not os.path.exists(content_file):
        print("❌ Content file not found. Please run gen_plan_content.py first.")
        return
    
    with open(content_file, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    # Generate presentation
    generator = PowerPointGenerator()
    generator.generate_presentation(data)
    generator.save_presentation()
    
    print("🎉 PowerPoint presentation generated successfully!")

if __name__ == "__main__":
    main()
